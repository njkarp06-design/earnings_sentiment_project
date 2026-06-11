"""
EarningsSentiment — v4 presentation
Narrative arc: Title → Problem → Pipeline → Scoring Engine → Architecture
              → Product Showcase (screenshots) → Product Views (screenshots)
              → Numbers → Stack → Roadmap → Next Step → Links → Career
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SHOTS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "screenshots")

def img(name):
    return os.path.join(SHOTS, name)

# ── Palette (used deliberately, not decoratively) ────────────────────────────
BG     = RGBColor(0x0D, 0x11, 0x17)   # near-black background
CARD   = RGBColor(0x1A, 0x23, 0x32)   # card / panel background
CARD2  = RGBColor(0x12, 0x1A, 0x26)   # darker card variant
TEAL   = RGBColor(0x00, 0xD4, 0x8A)   # primary accent — positive/success
BLUE   = RGBColor(0x38, 0xBD, 0xF8)   # secondary accent — informational
YELLOW = RGBColor(0xFB, 0xBF, 0x24)   # callout — warning / highlight
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)   # tertiary accent
RED    = RGBColor(0xF8, 0x71, 0x71)   # danger / emphasis (used sparingly)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xCB, 0xD5, 0xE1)   # light gray — readable body
GRAY   = RGBColor(0x94, 0xA3, 0xB8)   # medium gray — secondary text
DGRAY  = RGBColor(0x47, 0x55, 0x69)   # dark gray — metadata

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Primitive helpers ─────────────────────────────────────────────────────────

def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def box(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.line.width = 0
    return s


def t(slide, text, x, y, w, h,
      size=16, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def ml(slide, lines_cfg, x, y, w, h, align=PP_ALIGN.LEFT, spacing=2):
    """Multi-line textbox. lines_cfg = list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines_cfg):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i > 0:
            p.space_before = Pt(spacing)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def slide_title(slide, title, subtitle=None):
    bg(slide)
    box(slide, 0, Inches(0.08), W, Inches(0.05), TEAL)
    t(slide, title, Inches(0.55), Inches(0.25), Inches(12), Inches(0.72),
      size=38, bold=True, color=WHITE)
    if subtitle:
        t(slide, subtitle, Inches(0.55), Inches(1.02), Inches(12), Inches(0.45),
          size=16, color=TEAL, italic=True)


def stat_box(slide, number, label, x, y, w, h, num_color=TEAL):
    box(slide, x, y, w, h, CARD)
    box(slide, x, y, Inches(0.06), h, num_color)
    t(slide, number, x + Inches(0.2), y + Inches(0.12), w - Inches(0.3),
      Inches(h.inches * 0.55), size=42, bold=True, color=num_color)
    t(slide, label, x + Inches(0.2), y + Inches(h.inches * 0.55),
      w - Inches(0.3), Inches(h.inches * 0.42), size=13, color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# 1. TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, Inches(0.28), H, TEAL)
box(s, 0, 0, W, Inches(0.06), CARD2)

t(s, "//Developers_Institute_  Coding Bootcamp",
  Inches(0.55), Inches(0.18), Inches(9), Inches(0.42),
  size=13, color=DGRAY, italic=True)

t(s, "EarningsSentiment",
  Inches(0.55), Inches(0.85), Inches(10.5), Inches(1.5),
  size=78, bold=True, color=WHITE)

t(s, "Does how a CEO speaks on an earnings call predict what the stock does next?",
  Inches(0.55), Inches(2.5), Inches(9.8), Inches(0.9),
  size=22, color=TEAL, italic=True)

box(s, Inches(0.55), Inches(3.55), Inches(6.5), Inches(0.04), TEAL)

t(s, "Real-time AI platform  ·  earnings call NLP  ·  post-call price correlation",
  Inches(0.55), Inches(3.72), Inches(11), Inches(0.45),
  size=15, color=GRAY)

# 4 key stats
stats = [
    ("~10,000", "companies\nsearchable",   TEAL),
    ("5",       "AI signals\nper call",    BLUE),
    ("730",     "days of history\nper CEO", YELLOW),
    ("90%",     "token cost cut\nvia caching", PURPLE),
]
for i, (num, lbl, col) in enumerate(stats):
    bx = Inches(0.55 + i * 3.12)
    box(s, bx, Inches(4.35), Inches(2.95), Inches(1.7), CARD)
    box(s, bx, Inches(4.35), Inches(2.95), Inches(0.07), col)
    t(s, num, bx + Inches(0.2), Inches(4.45), Inches(2.6), Inches(0.85),
      size=44, bold=True, color=col)
    t(s, lbl, bx + Inches(0.2), Inches(5.3), Inches(2.6), Inches(0.6),
      size=13, color=GRAY)

t(s, "Natanel Karp  ·  2026",
  Inches(0.55), Inches(6.9), Inches(5), Inches(0.45),
  size=13, color=DGRAY)


# ════════════════════════════════════════════════════════════════════════════
# 2. THE PROBLEM  — minimal, punchy, emotional
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Faint giant background number for visual interest
t(s, "3,000", Inches(6.5), Inches(0.5), Inches(7), Inches(5.5),
  size=220, bold=True, color=RGBColor(0x18, 0x24, 0x30),
  align=PP_ALIGN.LEFT)

box(s, 0, 0, Inches(0.28), H, RED)

t(s, "THE PROBLEM",
  Inches(0.55), Inches(0.35), Inches(6), Inches(0.45),
  size=13, bold=True, color=RED)

ml(s, [
    ("Every quarter,",                  28, False, GRAY),
    ("~3,000 US companies hold earnings calls.",  36, True,  WHITE),
    ("Billions of dollars of trading",   28, False, GRAY),
    ("decisions follow.",                36, True,  WHITE),
    ("None of it is systematically",     28, False, GRAY),
    ("measured.",                        36, True,  WHITE),
], Inches(0.55), Inches(1.0), Inches(8.5), Inches(5.5), spacing=4)

box(s, Inches(0.55), Inches(6.62), Inches(8), Inches(0.04), TEAL)
t(s, "EarningsSentiment builds the measurement layer.",
  Inches(0.55), Inches(6.72), Inches(9), Inches(0.45),
  size=18, bold=True, color=TEAL)


# ════════════════════════════════════════════════════════════════════════════
# 3. HOW IT WORKS — simple 4-step pipeline
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "How It Works",
            "Four steps from filing to insight")

steps = [
    ("01", "INGEST",    BLUE,
     "SEC EDGAR & FMP\ntranscripts arrive\nwithin seconds\nof filing"),
    ("02", "SCORE",     TEAL,
     "Claude AI extracts\n5 signals from\nevery CEO's\nlanguage"),
    ("03", "CORRELATE", YELLOW,
     "1d, 3d, 7d stock\nreturns fetched\nautomatically\nvia yfinance"),
    ("04", "SURFACE",   PURPLE,
     "Feed, leaderboard,\ncompany pages,\ncalendar, and\nportfolio watchlist"),
]

BW = Inches(2.8)
BH = Inches(4.8)
GAP = Inches(0.62)
SX = Inches(0.55)

for i, (num, title, col, desc) in enumerate(steps):
    bx = SX + i * (BW + GAP)
    # card
    box(s, bx, Inches(1.6), BW, BH, CARD)
    # colour top strip
    box(s, bx, Inches(1.6), BW, Inches(0.12), col)
    # big step number (background)
    t(s, num, bx + Inches(0.1), Inches(1.72), BW - Inches(0.15), Inches(1.1),
      size=60, bold=True, color=RGBColor(0x22, 0x30, 0x45))
    # step title
    t(s, title, bx + Inches(0.22), Inches(2.7), BW - Inches(0.3), Inches(0.55),
      size=20, bold=True, color=col)
    # divider
    box(s, bx + Inches(0.22), Inches(3.3), BW - Inches(0.44), Inches(0.03), DGRAY)
    # description
    t(s, desc, bx + Inches(0.22), Inches(3.42), BW - Inches(0.35), Inches(2.5),
      size=16, color=LGRAY)
    # arrow (not after last)
    if i < len(steps) - 1:
        ax = bx + BW + Inches(0.12)
        t(s, "→", ax, Inches(3.8), GAP - Inches(0.12), Inches(0.55),
          size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Bottom callout
box(s, Inches(0.55), Inches(6.55), W - Inches(1.1), Inches(0.52), CARD2)
t(s, "Self-validating by design — every prediction is automatically back-tested against real post-call returns",
  Inches(0.75), Inches(6.63), W - Inches(1.4), Inches(0.38),
  size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# 4. THE SCORING ENGINE — the differentiator
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "The Scoring Engine",
            "What Claude extracts from every earnings call — the differentiating feature")

# Left: mock transcript card
TC_X = Inches(0.55)
TC_Y = Inches(1.6)
TC_W = Inches(5.2)
TC_H = Inches(5.6)

box(s, TC_X, TC_Y, TC_W, TC_H, CARD)
box(s, TC_X, TC_Y, TC_W, Inches(0.1), BLUE)
t(s, "TRANSCRIPT EXCERPT  —  AAPL Q1 2026",
  TC_X + Inches(0.2), TC_Y + Inches(0.18), TC_W - Inches(0.35), Inches(0.38),
  size=11, bold=True, color=BLUE)
box(s, TC_X + Inches(0.2), TC_Y + Inches(0.58), TC_W - Inches(0.4), Inches(0.03), DGRAY)
t(s, '"Our guidance for fiscal year 2026 reflects our\nconfidence in continued strong demand across\nour product lines. We are raising full-year EPS\nguidance to $7.20, up from our prior range.\n\nDuring Q&A, when pressed on supply constraints,\nmanagement responded: \'We have robust supply\nchain visibility and remain confident our\ncommitments will be met on schedule.\'"',
  TC_X + Inches(0.2), TC_Y + Inches(0.72), TC_W - Inches(0.35), Inches(3.1),
  size=14, color=LGRAY, italic=True)

box(s, TC_X + Inches(0.2), TC_Y + Inches(3.95), TC_W - Inches(0.4), Inches(0.03), DGRAY)
t(s, "Sent to Claude with cached system prompt  →",
  TC_X + Inches(0.2), TC_Y + Inches(4.08), TC_W - Inches(0.35), Inches(0.45),
  size=13, color=TEAL, bold=True)

# Big arrow in the middle
t(s, "→",
  Inches(5.9), Inches(3.7), Inches(0.8), Inches(0.65),
  size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Right: 5 output signal cards
signals = [
    ("Confidence Score",   "82 / 100",   TEAL,
     "Overall CEO language confidence — graded from crisis (0) to exceptional (100)"),
    ("Guidance Flag",      "RAISED ↑",    YELLOW,
     "Management raised, maintained, lowered, or withdrew forward guidance"),
    ("Q&A Defensiveness",  "2 / 10",     BLUE,
     "How evasive management became during analyst questions vs. prepared remarks"),
    ("Key Phrases",        "Top 3",      PURPLE,
     "The 3 verbatim phrases that moved the score most — direct from the transcript"),
    ("Trade Brief",        "2 sentences", GRAY,
     "Plain-English signal for a short-term trader: tone + most important fwd statement"),
]

SIG_X = Inches(6.85)
SIG_W = Inches(6.0)
SIG_H = Inches(0.95)
SIG_GAP = Inches(0.15)

for i, (name, value, col, desc) in enumerate(signals):
    sy = Inches(1.6) + i * (SIG_H + SIG_GAP)
    box(s, SIG_X, sy, SIG_W, SIG_H, CARD)
    box(s, SIG_X, sy, Inches(0.18), SIG_H, col)
    t(s, name,  SIG_X + Inches(0.32), sy + Inches(0.08),
      Inches(2.6), Inches(0.4), size=14, bold=True, color=WHITE)
    t(s, value, SIG_X + Inches(3.1), sy + Inches(0.1),
      Inches(2.7), Inches(0.4), size=18, bold=True, color=col, align=PP_ALIGN.RIGHT)
    t(s, desc,  SIG_X + Inches(0.32), sy + Inches(0.52),
      SIG_W - Inches(0.45), Inches(0.38), size=12, color=GRAY)

# Caching callout strip
STRIP_Y = Inches(1.6) + 5 * (SIG_H + SIG_GAP) + Inches(0.1)
box(s, SIG_X, STRIP_Y, SIG_W, Inches(0.58), CARD2)
t(s, "System prompt cached  →  ~80% latency reduction  ·  ~90% token-cost saving on cache hits",
  SIG_X + Inches(0.2), STRIP_Y + Inches(0.1), SIG_W - Inches(0.3), Inches(0.4),
  size=12, color=TEAL, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# 5. ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "Architecture",
            "Event-driven microservices — 5 containerised services, Kafka backbone")

svcs = [
    ("Ingestor",            "Python",      "SEC EDGAR\nFMP / Alpha Vantage\nPolls every 2 h\nseeds ~10k companies",        BLUE),
    ("Scoring\nService",    "Python",      "Consumes\nraw-transcripts\nClaude AI → 5 signals\n→ scored-transcripts",      TEAL),
    ("Correlation\nService","Python",      "Consumes\nscored-transcripts\nyfinance returns\n→ MongoDB",                    YELLOW),
    ("BFF\nREST API",       "Node.js",     "Express + JWT auth\n15+ endpoints\nReads MongoDB\nproduction-hardened",       PURPLE),
    ("Frontend",            "Next.js 14",  "6 product views\nRecharts charts\nSSE streaming\nTailwind CSS",               RED),
]

BW2 = Inches(2.35)
BH2 = Inches(3.95)
GAP2 = Inches(0.16)
SX2  = Inches(0.42)

for i, (name, lang, desc, col) in enumerate(svcs):
    bx = SX2 + i * (BW2 + GAP2)
    box(s, bx, Inches(1.6), BW2, BH2, CARD)
    box(s, bx, Inches(1.6), BW2, Inches(0.1), col)
    t(s, name,  bx + Inches(0.15), Inches(1.75), BW2 - Inches(0.25), Inches(0.72),
      size=15, bold=True, color=WHITE)
    t(s, lang,  bx + Inches(0.15), Inches(2.5),  BW2 - Inches(0.25), Inches(0.38),
      size=12, bold=True, color=col)
    box(s, bx + Inches(0.15), Inches(2.9), BW2 - Inches(0.3), Inches(0.03), DGRAY)
    t(s, desc,  bx + Inches(0.15), Inches(3.0),  BW2 - Inches(0.25), Inches(2.2),
      size=13, color=LGRAY)
    if i < len(svcs) - 1:
        ax = bx + BW2 + Inches(0.01)
        t(s, "→", ax, Inches(3.3), GAP2 + Inches(0.06), Inches(0.55),
          size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Kafka row
box(s, Inches(0.42), Inches(5.72), Inches(9.6), Inches(0.52), CARD2)
t(s, "Apache Kafka   |   raw-transcripts  ·  raw-prices  ·  scored-transcripts",
  Inches(0.62), Inches(5.8), Inches(9.2), Inches(0.36),
  size=13, bold=True, color=TEAL)

# MongoDB
box(s, Inches(10.15), Inches(5.72), Inches(2.76), Inches(0.52), CARD2)
t(s, "MongoDB  |  price_reactions · scores · users",
  Inches(10.28), Inches(5.8), Inches(2.55), Inches(0.36),
  size=11, bold=True, color=YELLOW)

# AWS banner
box(s, Inches(0.42), Inches(6.38), W - Inches(0.84), Inches(0.68), RGBColor(0x0E, 0x17, 0x24))
t(s, "Production  :  AWS ECS Fargate  ·  Amazon MSK  ·  MongoDB Atlas  ·  ECR  ·  ALB x2  ·  Terraform  ·  GitHub Actions CI/CD",
  Inches(0.65), Inches(6.5), W - Inches(1.2), Inches(0.42),
  size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# 6. PRODUCT SHOWCASE — Feed + Company page (hero screenshots)
# Screenshots sit within y=1.55 to y=6.82; caption strip y=6.85 to y=7.22
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "The Product — Live",
            "Real data. Real companies. Running now.")

SH6 = Inches(5.27)   # screenshot height — leaves room for caption below

# Feed — left two-thirds
s.shapes.add_picture(img("feed.png"),    Inches(0.42), Inches(1.55), Inches(7.62), SH6)
box(s, Inches(0.42), Inches(6.85), Inches(7.62), Inches(0.36), CARD2)
t(s, "Earnings Feed  —  live CEO confidence scores, key phrases, 1d / 3d / 7d returns",
  Inches(0.58), Inches(6.9), Inches(7.35), Inches(0.28), size=13, color=TEAL, bold=True)

# Company page — right third
s.shapes.add_picture(img("company.png"), Inches(8.22), Inches(1.55), Inches(4.69), SH6)
box(s, Inches(8.22), Inches(6.85), Inches(4.69), Inches(0.36), CARD2)
t(s, "Company deep-dive  —  SNOW: 12 calls tracked, +35.95% last 7D",
  Inches(8.35), Inches(6.9), Inches(4.45), Inches(0.28), size=13, color=BLUE, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# 7. PRODUCT VIEWS — Leaderboard + Sectors + Portfolio + Inspect
# 2x2 grid: SH=2.55, label=0.30, gap=0.14 — fits within 7.5" height
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "The Product — All Views",
            "Six views built, all running on AWS")

_SW  = Inches(6.48)   # screenshot width  (2 cols + 0.19" gap = 13.15" < 13.33")
_SH  = Inches(2.55)   # screenshot height
_LH  = Inches(0.30)   # label height
_GX  = Inches(0.19)   # horizontal gap between columns
_GY  = Inches(0.14)   # vertical gap between rows
_SX0 = Inches(0.42)   # left edge
_SY0 = Inches(1.55)   # top edge

shot_grid = [
    ("leaderboard.png", "Leaderboard  —  CRWD 91% win rate · ranked by avg 7D return", TEAL,   0, 0),
    ("sectors.png",     "Sector Pulse  —  Energy +3.77% · Technology 199 calls",        BLUE,   1, 0),
    ("portfolio.png",   "Portfolio  —  DELL +32.76% 1D, AAPL +7.96% 7D",               PURPLE, 0, 1),
    ("inspect.png",     "Call Inspect  —  AMD post-earnings drift · 11 prior calls",    YELLOW, 1, 1),
]

for (fname, label, col, ci, ri) in shot_grid:
    sx = _SX0 + ci * (_SW + _GX)
    sy = _SY0 + ri * (_SH + _LH + _GY)
    s.shapes.add_picture(img(fname), sx, sy, _SW, _SH)
    box(s, sx, sy + _SH, _SW, _LH, CARD2)
    t(s, label, sx + Inches(0.12), sy + _SH + Inches(0.05),
      _SW - Inches(0.2), _LH - Inches(0.06), size=12, color=col, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# 7. BY THE NUMBERS — stats wall (most impactful slide)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "By the Numbers",
            "The scale of what was built")

numbers = [
    ("~10,000", "companies\nsearchable from EDGAR",            TEAL),
    ("5",       "AI signals extracted\nfrom every call",       BLUE),
    ("199",     "Technology sector calls\nacross 17 companies", YELLOW),
    ("90%",     "token-cost reduction\nvia prompt caching",    PURPLE),
    ("91%",     "CRWD win rate —\nhighest on leaderboard",     TEAL),
    ("5",       "microservices\ndeployed on AWS",              BLUE),
    ("+35.95%", "SNOW last 7D return\nafter high-score call",  YELLOW),
    ("15+",     "REST API endpoints\nwith JWT auth",           PURPLE),
]

NW = Inches(2.85)
NH = Inches(2.38)
NGAP_X = Inches(0.42)
NGAP_Y = Inches(0.3)
NSX = Inches(0.68)
NSY = Inches(1.6)

for i, (num, lbl, col) in enumerate(numbers):
    row = i // 4
    ci  = i % 4
    nx = NSX + ci * (NW + NGAP_X)
    ny = NSY + row * (NH + NGAP_Y)
    box(s, nx, ny, NW, NH, CARD)
    box(s, nx, ny, NW, Inches(0.1), col)
    t(s, num, nx + Inches(0.2), ny + Inches(0.18), NW - Inches(0.3),
      Inches(1.18), size=52, bold=True, color=col)
    t(s, lbl, nx + Inches(0.2), ny + Inches(1.42), NW - Inches(0.3),
      Inches(0.85), size=14, color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# 8. STACK — grouped by layer, with rationale for bold choices
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "Tech Stack", "Tools, languages, and services — grouped by layer")

groups = [
    ("Data & AI Layer", TEAL, [
        "Python 3.12",
        "Apache Kafka  (3 topics)",
        "SEC EDGAR API",
        "FMP / Alpha Vantage",
        "Claude claude-sonnet-4-6",
        "Anthropic SDK + prompt caching",
        "yfinance (price data)",
        "APScheduler  (ingest schedule)",
    ]),
    ("Application Layer", BLUE, [
        "Node.js  /  Express",
        "MongoDB  +  Mongoose",
        "JWT authentication",
        "bcryptjs  (password hashing)",
        "Resend  (email alerts)",
        "REST API  (15+ endpoints)",
        "Flask sidecar  (ingest trigger)",
        "Docker  /  Docker Compose",
    ]),
    ("Frontend & Infra Layer", YELLOW, [
        "Next.js 14  (App Router)",
        "React 18  +  Tailwind CSS",
        "Recharts  (data visualisation)",
        "AWS ECS Fargate  (5 services)",
        "Amazon MSK  (managed Kafka)",
        "MongoDB Atlas  (cloud DB)",
        "ECR  +  ALB x2",
        "Terraform  +  GitHub Actions",
    ]),
]

GW = Inches(3.95)
GH = Inches(5.55)
GSX = Inches(0.55)

for i, (title, col, items) in enumerate(groups):
    gx = GSX + i * (GW + Inches(0.42))
    box(s, gx, Inches(1.58), GW, GH, CARD)
    box(s, gx, Inches(1.58), GW, Inches(0.12), col)
    t(s, title, gx + Inches(0.2), Inches(1.75), GW - Inches(0.3), Inches(0.5),
      size=17, bold=True, color=col)
    box(s, gx + Inches(0.2), Inches(2.28), GW - Inches(0.4), Inches(0.04), DGRAY)
    for j, item in enumerate(items):
        t(s, "  " + item, gx + Inches(0.2), Inches(2.42 + j * 0.59),
          GW - Inches(0.3), Inches(0.55), size=14, color=LGRAY)

# Key decision callouts
box(s, Inches(0.55), Inches(7.1), W - Inches(1.1), Inches(0.28), CARD2)
t(s, "Kafka chosen for decoupled stream processing  ·  Prompt caching chosen for cost/latency  ·  Terraform for reproducible AWS infra",
  Inches(0.75), Inches(7.13), W - Inches(1.4), Inches(0.22),
  size=12, color=DGRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# 9. ROADMAP — Gantt chart
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "Roadmap", "Next 4 priorities — ordered by impact")

ROAD_ITEMS = [
    ("1  FMP Data Depth",       TEAL,
     "Upgrade tier  ·  forced re-backfill endpoint  ·  full 5-year history per ticker",
     0.0, 3.0),
    ("2  Predictability Score", BLUE,
     "Score-to-return scatter  ·  Pearson r regression  ·  CEO Predictability badge",
     2.5, 6.0),
    ("3  UX Clarity",           YELLOW,
     "Tooltips for % win, Avg Score, thresholds, and pending return badges",
     5.5, 8.2),
    ("4  Score Trend + Sector", PURPLE,
     "Confidence trend over 4-6 quarters  ·  returns vs. sector benchmark",
     7.5, 10.8),
]

CL = Inches(0.45)
LW = Inches(3.5)
CW2 = W - Inches(0.9)
CT = Inches(1.6)
RH = Inches(1.25)
RGAP = Inches(0.12)
TWKS = 11.0
BAR_AREA = CW2 - LW

def wx(w):
    return CL + LW + (w / TWKS) * BAR_AREA

for i, (title, col, desc, start, end) in enumerate(ROAD_ITEMS):
    ry = CT + i * (RH + RGAP)
    row_fill = CARD if i % 2 == 0 else CARD2
    box(s, CL, ry, CW2, RH, row_fill)
    t(s, title, CL + Inches(0.18), ry + Inches(0.1),
      LW - Inches(0.25), Inches(0.48), size=15, bold=True, color=col)
    t(s, desc,  CL + Inches(0.18), ry + Inches(0.6),
      LW - Inches(0.25), Inches(0.58), size=12, color=GRAY)
    bx = wx(start)
    bw = wx(end) - bx
    box(s, bx, ry + Inches(0.22), bw, Inches(0.82), col)
    wk = f"Wk {int(start)+1}-{int(end)}"
    t(s, wk, bx + Inches(0.1), ry + Inches(0.38),
      bw - Inches(0.1), Inches(0.42), size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)

# Column separator
box(s, CL + LW, CT, Inches(0.03), len(ROAD_ITEMS)*(RH + RGAP), TEAL)

# Week markers
for w in range(0, 12, 2):
    mx = wx(w)
    box(s, mx, CT, Inches(0.02), len(ROAD_ITEMS)*(RH + RGAP), RGBColor(0x24, 0x33, 0x4A))
    t(s, f"W{w}", mx - Inches(0.18), CT + len(ROAD_ITEMS)*(RH + RGAP) + Inches(0.06),
      Inches(0.4), Inches(0.3), size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# 10. MY NEXT STEP
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "My Next Step")

next_steps = [
    ("1", "Fix FMP Data Depth",
     "Upgrade FMP tier and ship a forced re-backfill endpoint. "
     "Data quality is the foundation — without a full 5-year call history the predictability model is built on sand.",
     TEAL),
    ("2", "Ship the CEO Predictability Score",
     "Build the per-company score-to-return scatter with a Pearson r regression line and surface the "
     "Predictability Score prominently. This is the single most differentiated feature on the platform.",
     BLUE),
    ("3", "UX Clarity",
     "Add inline tooltips for every metric — % win, Avg Score, score thresholds, pending return badges. "
     "Low-lift, high-impact: makes the platform accessible to non-quant users without dumbing it down.",
     YELLOW),
    ("4", "Score Trend + Sector-Relative Returns",
     "Show confidence score trend over the last 4-6 quarters and post-call returns adjusted for sector "
     "movement. Adds the analytical depth expected by sophisticated users and portfolio managers.",
     PURPLE),
]

for i, (num, title, body, col) in enumerate(next_steps):
    sy = Inches(1.55) + i * Inches(1.38)
    box(s, Inches(0.45), sy, W - Inches(0.9), Inches(1.25), CARD)
    box(s, Inches(0.45), sy, Inches(0.7), Inches(1.25), col)
    t(s, num, Inches(0.45), sy + Inches(0.32),
      Inches(0.7), Inches(0.58), size=26, bold=True, color=BG, align=PP_ALIGN.CENTER)
    t(s, title, Inches(1.3), sy + Inches(0.1),
      Inches(3.8), Inches(0.5), size=18, bold=True, color=col)
    t(s, body, Inches(1.3), sy + Inches(0.58),
      W - Inches(1.75), Inches(0.62), size=14, color=LGRAY)


# ════════════════════════════════════════════════════════════════════════════
# 11. LINKS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "Links")

links = [
    ("GitHub",          "github.com/njkarp06-design/earnings_sentiment_project",  TEAL),
    ("Demo Video",      "loom.com/share/f538f79d328f420aac248f0383ebb3ee",          BLUE),
    ("Deployed App",    "njkarp06-design.github.io/earnings_sentiment_project",    YELLOW),
]

for i, (label, url, col) in enumerate(links):
    ry = Inches(1.6) + i * Inches(1.05)
    box(s, Inches(1.5), ry, W - Inches(3.0), Inches(0.88), CARD)
    box(s, Inches(1.5), ry, Inches(0.2), Inches(0.88), col)
    t(s, label, Inches(1.85), ry + Inches(0.13), Inches(3.0), Inches(0.45),
      size=18, bold=True, color=col)
    t(s, url,   Inches(5.1),  ry + Inches(0.15), Inches(7.5), Inches(0.58),
      size=16, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# 12. CAREER PREP
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "Career Prep")

# Name + bio block
box(s, Inches(0.45), Inches(1.55), W - Inches(0.9), Inches(1.35), CARD)
box(s, Inches(0.45), Inches(1.55), Inches(0.22), Inches(1.35), TEAL)
t(s, "Natanel Karp",
  Inches(0.82), Inches(1.65), Inches(6), Inches(0.6),
  size=28, bold=True, color=WHITE)
t(s, "Junior full-stack developer who treats AI tooling as part of the stack.  "
     "Built this project end-to-end.  Hackathon Winner 2026.",
  Inches(0.82), Inches(2.22), Inches(11.5), Inches(0.55),
  size=14, color=LGRAY)

# 4 key highlights from CV
highlights = [
    ("Full-Stack Dev Intern  —  Horizon Trade Tech (May-Jun 2026)",
     "In-memory LRU cache: ~700x speedup (97ms to 0.14ms p50). Solved 502 errors traced to cold EBS snapshot reads on Fargate.",
     BLUE),
    ("AI Data Analyst Intern  —  Elevata Ventures / Brevan Howard",
     "Built AI-driven investment evaluation system with Claude Code. Applied to 45+ startups. Briefed Brevan Howard on AI strategy.",
     PURPLE),
    ("Education  —  Developers Institute + Hasmonean A-Levels",
     "AI & Full-Stack Diploma (2026). A-Levels: Mathematics 96%, Further Mathematics, Economics, Psychology.",
     YELLOW),
    ("Achievements",
     "Author: 'Rumination' (60,000-word manuscript).  Gold, UKMT Senior Maths Challenge 2024.  Competitive distance runner.",
     TEAL),
]

for i, (title, body, col) in enumerate(highlights):
    row = i // 2
    ci2 = i % 2
    hx = Inches(0.45 + ci2 * 6.5)
    hy = Inches(3.1) + row * Inches(1.38)
    box(s, hx, hy, Inches(6.1), Inches(1.2), CARD)
    box(s, hx, hy, Inches(0.18), Inches(1.2), col)
    t(s, title, hx + Inches(0.32), hy + Inches(0.1),
      Inches(5.65), Inches(0.42), size=13, bold=True, color=col)
    t(s, body,  hx + Inches(0.32), hy + Inches(0.55),
      Inches(5.65), Inches(0.6), size=12, color=GRAY)

# Links row
link_items = [("CV",  "drive.google.com/file/d/1xvYOBGgMApOvlOSh_a2i1uAQ64rfpEqE/view", TEAL), ("LinkedIn", "linkedin.com/in/natanel-karp-6890aa3b3", BLUE)]
LRW = (W - Inches(0.9)) / 3
for i, (lbl, url, col) in enumerate(link_items):
    lx = Inches(0.45) + i * LRW
    box(s, lx, Inches(5.95), LRW - Inches(0.12), Inches(0.85), CARD2)
    t(s, lbl, lx + Inches(0.2), Inches(6.02), LRW - Inches(0.5), Inches(0.35),
      size=15, bold=True, color=col)
    t(s, url,  lx + Inches(0.2), Inches(6.4), LRW - Inches(0.5), Inches(0.35),
      size=13, color=WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "EarningsSentiment_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
